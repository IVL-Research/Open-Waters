import os
import datetime
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objs as go
import matplotlib.pyplot as plt
import glob
import pptx
from pptx.util import Inches, Pt
import ipywidgets as widgets
from ipywidgets import interact, Layout

class TimeSeriesConstrictor:
    """
    A class to contain a pandas dataframe along with functions to modify the data in it.
    If the class methods are used, the original data in the dataframe is always kept as it is.
    Instead, new columns are created with the modified data along with how they were created.
    """

    def __init__(self, dataframe=None, metadata=None, description=None):
        if isinstance(dataframe, type(pd.DataFrame())):
            self.dataframe = dataframe
        if isinstance(metadata, type(dict())):
            self.metadata = metadata
        else:
            self.metadata = dict()
        if isinstance(description, type(dict())):
            self.description = description
        else:
            self.description = dict()

    def plot(self, y_column="all", default_size=3, **kwargs):
        """
        Returns interactive plot.
        y_column is the column name in self.dataframe,
        can be either a string or a list of strings.

        """
        fig = go.Figure()

        if y_column == "all":
            y_column = [str(col) for col in self.dataframe.columns]
        elif not isinstance(y_column, type([])):
            y_column = [y_column]
        for col in y_column:
            marker_dict = dict()
            temp_data = self.dataframe[col]
            mode = "lines+markers"
            marker_dict["size"] = default_size
            line_dict = {'width': 1}
            try:
                mode = self.metadata[col]["plot_mode"]
                marker_dict["symbol"] = self.metadata[col]["plot_markers"]
                marker_dict["size"] = self.metadata[col]["marker_size"]
            except:
                pass
            fig.add_trace(
                go.Scattergl(
                    x=temp_data.index,
                    y=temp_data,
                    name=col,
                    mode=mode,
                    marker=marker_dict,
                    line=line_dict,
                    connectgaps=False
                )
            )
        fig.update_layout(showlegend=True)

        return fig

    def plot_static(self, y_column, save_name='', **kwargs):
        """
        Returns and/or saves static plot, for export of results.
        y_column is the column name in self.dataframe.
        """

        # Check if outlier column
        if 'outlier' in y_column:
            # Plot outliers as points above used data column
            if y_column not in self.metadata.keys():
                return
            temp_df = self.dataframe[[y_column]].copy()
            temp_df['Time'] = temp_df.index
            fig, ax = plt.subplots(figsize=(12, 5))
            self.dataframe[self.metadata[y_column]['used_data_column']].plot(ax=ax)
            temp_df.plot(x='Time', y=y_column, kind='scatter', color='DarkOrange', ax=ax)
        else:
            # Set fig size suitable for time series
            plt.figure(figsize=(12, 5))
            # Plot y_column
            self.dataframe[y_column].plot()

        # Set axis labels
        try:
            plt.ylabel(self.description[y_column]['unit'])
        except KeyError:
            pass

        # Set title and grid options
        plt.title(y_column)
        plt.grid()
        plt.xticks(rotation=90)

        if not save_name == '':
            plt.savefig(save_name, bbox_inches='tight')
            plt.close()
        else:
            plt.show()

    def parameter_tuning(self, methods_dict, start_date='2020-01-01', stop_date='2021-01-01'):
        """
        Interactive plot to tune pre-processing method parameters using sliders.
        The plot is updated upon changing slider values, target column or method.
        The method does not store the calculated values, and the values are calculated upon each update.

        The methods and their respective options and settings to be evaluated are defined in the methods_dict.
        E.g:

        methods_dict = {'outlier_detection':{'median_lim':{'min':0, 'max':5, 'step':0.05,
                                                        'description': '"Median limit"'},
                                        'window_size':{'min':0, 'max':30, 'step':1,
                                                        'description': '"Window size"'}},
                    'find_frozen_values':{'var_lim_low':{'min':0, 'max':5, 'step':0.05,
                                                        'description': '"Var lim low"'},
                                        'window_size':{'min':0, 'max':30, 'step':1,
                                                        'description': '"Window size"'}},
                   'out_of_range_detection':{'min_limit':{'min':-100, 'max':100, 'step':5,
                                                        'description': '"Min limit"'},
                                        'max_limit':{'min':0, 'max':300, 'step':5,
                                                        'description': '"Max limit"'}}}

        """

        # Date picker widgets
        start_date = widgets.DatePicker(description='Start', disabled=False, value=pd.to_datetime(start_date))
        display(start_date)
        stop_date = widgets.DatePicker(description='Stop', disabled=False, value=pd.to_datetime(stop_date))
        display(stop_date)

        # Default parameter sliders
        global s1, s2, s3, s4
        s1 = widgets.FloatSlider(min=0, max=100, step=1, value=1,
                                 description='', layout=Layout(width='50%'))
        s2 = widgets.FloatSlider(min=0, max=100, step=1, value=1,
                                 description='', layout=Layout(width='50%'))
        s3 = widgets.FloatSlider(min=0, max=100, step=1, value=1,
                                 description='', layout=Layout(width='50%'))
        s4 = widgets.FloatSlider(min=0, max=100, step=1, value=1,
                                 description='', layout=Layout(width='50%'))

        # Method and target column selection
        method_list = widgets.Select(
            options=list(methods_dict.keys()),
            value=list(methods_dict.keys())[0],
            description='Methods',
            disabled=False, layout=Layout(width='50%', height='80px'))

        column_list = widgets.Select(
            options=self.dataframe.columns,
            value=self.dataframe.columns[0],
            description='Variables',
            disabled=False, layout=Layout(width='50%', height='200px'))

        # Initialize figure
        fig = go.Figure()
        fig.update_layout(legend=dict(orientation="h",
                                      yanchor="bottom",
                                      y=1.02,
                                      xanchor="right",
                                      x=1),
                          margin=dict(l=20, r=20, t=25, b=20))

        @interact
        def update_plot(method=method_list, column=column_list, default_1=s1, default_2=s2, default_3=s3, default_4=s4):

            """
            Each time the update_plot function is called, the slider settings are updated
            and a new calculation string is created.
            """

            count = 1
            calc_str = 'self.' + method + '(target_column="' + column + '"'  # Initialize calculation string

            for slider in methods_dict[method]:

                # Update slider settings
                for option in methods_dict[method][slider]:
                    option_str = 's' + str(count) + '.' + option + '=' + str(methods_dict[method][slider][option])
                    exec(option_str)

                # Add parameter options to calculation string
                value = eval('default_' + str(count))
                if slider == 'window_size':
                    value = int(value)
                calc_str = calc_str + ',' + slider + '=' + str(value)

                count += 1

            # Disable unused sliders
            while count <= 4:
                option_str = 's' + str(count) + '.' + 'disabled' + '=True'
                exec(option_str)
                count += 1

            calc_str = calc_str + ',return_only=True)'

            # Calculate results
            df = eval(calc_str)  # Evaluate the calculation string
            df = df.loc[start_date.value:stop_date.value]  # Preprocessed results

            plot_df = pd.DataFrame(self.dataframe[column].loc[start_date.value:stop_date.value])  # Raw data
            plot_df[method] = df
            plot_x = plot_df.index

            # Plot results
            fig.data = []  # Clear figure data

            # Methods to plot only markers and not lines
            marker_functions = ["outlier_detection"]
            line_dict = {}
            marker_dict = {}

            for col in plot_df.columns:
                temp_data = plot_df[col]
                if col in marker_functions:
                    mode = "markers"
                    marker_dict["symbol"] = "circle-open"
                    marker_dict["size"] = 6
                else:
                    mode = "lines+markers"
                    marker_dict["size"]= 3
                    line_dict['width']= 1
                fig.add_trace(
                    go.Scattergl(  # Using Scattergl instead of Scatter to speed up rendering
                        x=plot_x,
                        y=temp_data,
                        name=col,
                        mode=mode,
                        marker=marker_dict,
                        line=line_dict,
                        connectgaps=False
                    ))

            fig.show()

    def create_metadata(self, metadata_dict, target_column):

        self.metadata[target_column] = dict()
        for key in metadata_dict.keys():
            self.metadata[target_column][key] = metadata_dict[key]

    def create_description(self, target_column, var_def=None):

        self.description[target_column] = dict()

        # Read from variable definition file
        if var_def:
            try:
                self.description[target_column]['info'] = var_def['Description'][target_column]
                self.description[target_column]['unit'] = var_def['EngineeringUnit'][target_column]
                self.description[target_column]['min limit'] = var_def['Min limit'][target_column]
                self.description[target_column]['max limit'] = var_def['Max limit'][target_column]
                self.description[target_column]['max constant values'] = var_def['MaxConstantValues'][target_column]
            except KeyError:
                pass

        # Calculate descriptive statistics
        self.description[target_column]["NaN count"] =  self.dataframe[target_column].isnull().sum()
        statistics = self.dataframe[target_column].describe()
        for ix, stat in enumerate(statistics):
            self.description[target_column][statistics.index[ix]] = stat

    def add_preprocessed_column(self, data, method, parameters, target_column, column_name='preprocessed', metadata={},
                                **kwargs):
        """
        Add data preprocessed by arbitrary method to TSC dataframe and store the corresponding metadata.

        Parameters
        -----------
        data : pd.Series or np.array
            Preprocessed data, must be same length as TSC.dataframe
        method : str
            Name of method used for preprocessing (including package name), eg 'scipy.some_method'
        parameters : dict
            Dict of parameter values used for preprocessing, eg {'param1': 3, 'param2': 0.14}
        used_column : str
            Name of column used for preprocessing
        column_name : str
            Name of added column
        metadata : dict
            Metadata returned by the preprocessing method and/or plot settings
        """

        column_name = self.create_target_column(column_name)
        self.dataframe[column_name] = data
        md_dict = {'method': method, 'target_column': target_column}

        for item in parameters:
            md_dict[item] = parameters[item]

        for item in metadata:
            md_dict[item] = metadata[item]

        self.create_metadata(md_dict, column_name)

        self.create_description(column_name)

    def one_hot_encoder(self, target_column, **kwargs):
        """
        Perform One Hot Encoding using pandas.get_dummies() and store data and metadata to the TSC object.

        Parameters
        -----------
        target_column : str
            Target column in TSC.dataframe
        **kwargs : dict
            Optional keyword arguments to pd.get_dummies()
        """
        encoded = pd.get_dummies(self.dataframe[target_column], **kwargs)

        for col in encoded:
            new_column = self.create_target_column(col)
            add_preprocessed_column(self=self,
                                    data=encoded[col],
                                    method='one_hot_encoded',
                                    parameters=kwargs,
                                    used_column=target_column,
                                    column_name=new_column,
                                    metadata={})

    def smoothing(
        self,
        target_column,
        smoothing_technique="exponential",
        alpha=0.2,
        window_size=3,
        output_column_name="preprocessed",
        **kwargs
    ):
        """
        Smoothens data.
        """
        # Create name for column
        new_column = self.create_target_column(output_column_name)

        if smoothing_technique == "exponential":
            
            # Create own column with smoothened data
            # Är det enl nedan jag behöver skriva för att skicka vidare värden på alpha? Trodde det skulle räcka att bara
            # skriva "alpha"?
            self.dataframe[new_column] = self.dataframe[
                target_column].ewm(alpha=alpha, adjust=False).mean()
            
            self.dataframe[new_column][self.dataframe[target_column].isnull()]=np.nan
            
            # Create metadata dictionary for column with smoothened data
            metadata_dict = {"method": "smoothing",
                         "used_data_column": target_column,
                         "smoothing_technique": smoothing_technique,
                         "alpha": alpha}
        
        elif smoothing_technique == "movAv":
            self.dataframe[new_column] = self.dataframe[target_column].rolling(window=window_size).mean()
             
            # Create metadata dictionary for column with smoothened data
            metadata_dict = {"method": "smoothing",
                         "used_data_column": target_column,
                         "smoothing_technique": smoothing_technique,
                         "window_size": window_size}

        self.create_metadata(metadata_dict, new_column)
        
        self.create_description(new_column)

    def outlier_detection(
        self,
        target_column,
        outlier_dist="robust",
        window_size=3,
        median_lim=1,
        mode="run",
        output_column_name="preprocessed",
        outlier_column_name="outliers",
        return_only=False,
        **kwargs
    ):
        """
        Find outliers.
        Creates one dataframe column called "outliers" where outliers are
        kept and one dataframe column called "preprocessed" where outliers are removed.
        """

        # Constants should always be stated at the top and explained
        CONSTANT_1 = (
            0.675  # Standard distribution constant, portion of data within 1 STD.
        )

        # Create temporary dataframe to use only inside this method
        outlier_detection_temp_df = pd.DataFrame()

        # Median of previous (window_size) values
        val = (
            self.dataframe[target_column]
            .rolling(window_size)
            .median(**kwargs)
            .shift(1, axis=0)
        )

        if outlier_dist == "robust":
            # The robust method compares the deviation of a data point (target_column) and the median of the previous window (val) 
            # with the MAD of the entire data series. 
            # MAD=Median absolute deviation=The median of the absolute deviations from the datas median
            
            # MATLAB:
            # madVal=median(abs(dataNoNans-median(dataNoNans)));
            
            # .mad in pandas is mean abs. deviation
            #mad_val = self.dataframe[target_column].mad(**kwargs)
            
            # New attempt:
            mad_val=(self.dataframe[target_column]-self.dataframe[target_column].median()).abs().median()

            target_median = self.dataframe[target_column].median()
            abs_diff_to_median = (self.dataframe[target_column] - target_median).abs()
            median_of_abs_diff_to_median = abs_diff_to_median.median()


            # Compare value with val
            outlier_detection_temp_df["Test"] = CONSTANT_1 * (
                (self.dataframe[target_column] - val).abs() / mad_val
            )

            
        elif outlier_dist == "gaussian":
            # Standard deviation
            dataStd = self.dataframe[target_column].std(**kwargs)

            # Compare value with val
            outlier_detection_temp_df["Test"] = (
                self.dataframe[target_column] - val
            ).abs() / dataStd
        else:
            print("No valid outlier_dist specified, doing nothing")
            outlier_detection_temp_df["Test"] = np.nan
            self.metadata[new_column]["outlier_dist"] = "None"

        if mode == "run":
            # Create column where 0="not outlier", 1="outlier"
            outlier_detection_temp_df["anomalyVec"] = np.zeros(len(self.dataframe))
            outlier_detection_temp_df.loc[
                outlier_detection_temp_df["Test"] > median_lim, "anomalyVec"
            ] = 1

            if not return_only:

                # Create name for column
                new_column = self.create_target_column(outlier_column_name)
                new_column_2 = self.create_target_column(output_column_name)

                # Create metadata dictionary for outliers
                metadata_dict = {"method": "outlier_detection",
                                 "used_data_column": target_column,
                                 "outlier_dist": outlier_dist,
                                 "window_size": window_size,
                                 "median_lim": median_lim,
                                 "mode": mode,
                                 "plot_mode": "markers",
                                 "plot_markers": "circle-open",
                                 "marker_size": 8}
                self.create_metadata(metadata_dict, new_column)

                # Create metadata dictionary for column where outliers have been removed

                metadata_dict = {"method": "outlier_detection",
                                 "used_data_column": target_column,
                                 "outlier_dist": outlier_dist,
                                 "window_size": window_size,
                                 "median_lim": median_lim,
                                 "mode": mode}
                self.create_metadata(metadata_dict, new_column_2)

                # Create own column with outliers
                self.dataframe[new_column] = self.dataframe[
                    target_column
                ] * outlier_detection_temp_df["anomalyVec"].replace(0, np.nan)

                # The number of data points that were not possible to check with current test due to Nans inside the moving window.
                # Nans in original data are exclude from the sum.
                self.metadata[new_column]["NonCheckedData"] = dict()
                self.metadata[new_column]["NonCheckedData"] = (
                    outlier_detection_temp_df["Test"].isna().sum()
                    - self.dataframe[target_column].isna().sum()
                )

                self.metadata[new_column]["AnomalyProportion"] = (
                    outlier_detection_temp_df["anomalyVec"].sum()
                    / outlier_detection_temp_df["anomalyVec"].count()
                )

                # Create column where outliers have been removed from the target data
                self.dataframe[new_column_2] = self.dataframe[target_column].loc[
                    outlier_detection_temp_df["anomalyVec"] < 0.5
                ]

                # Add description for outliers and pre-processing columns
                self.create_description(new_column)
                self.description[new_column]["info"] = (
                    "Data where all NON-outliers of " + target_column + " are set to nan"
                )
                self.description[new_column][
                    "NonCheckedData"
                ] = "Number of data points that were not possible to \
                evaluate with the current test due to Nans inside the moving window. Nans in original data are exclude from the sum."

                self.create_description(new_column_2)
                self.description[new_column_2]["info"] = (
                        "Data where outliers of " + target_column + " are set to nan"
                )
            else:
                outlier_return = self.dataframe[target_column] * \
                                 outlier_detection_temp_df["anomalyVec"].replace(0, np.nan)
                return outlier_return

    def resample(self, target_column, **kwargs):
        """Implementation of pandas resample function with added functionality to store resampling metadata
        Parameters
        -----------
        target_column : str
            Target column in TSC.dataframe
        **kwargs : dict
            Optional keyword arguments to pd.resample()
            """
        resampled_data = self.dataframe[target_column].resample(**kwargs)
        add_preprocessed_column(self=self,
                                data=resampled_data,
                                method='resample',
                                parameters=kwargs,
                                used_column=target_column,
                                column_name='resampled',
                                metadata={})

    def interpolate(self, target_column, **kwargs):
        """Implementation of pandas interpolate function with added functionality to store interpolation metadata
        Parameters
        -----------
        target_column : str
            Target column in TSC.dataframe
        **kwargs : dict
            Optional keyword arguments to pd.interpolate()
            """
        interpolated_data = self.dataframe[target_column].interpolate(**kwargs)
        add_preprocessed_column(self=self,
                                data=interpolated_data,
                                method='interpolate',
                                parameters=kwargs,
                                used_column=target_column,
                                column_name='interpolated',
                                metadata={})

    def read_excel(self, path, index_col=0, start_date='1970-01-01', stop_date='2100-01-01', **kwargs):
        self.dataframe = pd.read_excel(path, index_col=index_col, engine="openpyxl",**kwargs)
        self.dataframe.index = pd.to_datetime(self.dataframe.index)
        if (self.dataframe.index.dtype != self.dataframe.index.tz_localize(None).dtype):
            print("Note that timezone information has been removed from index at import!")
        self.dataframe.index = self.dataframe.index.tz_localize(None)
        self.dataframe = self.dataframe[start_date:stop_date]
        
        for column in self.dataframe.columns:
            self.create_description(column)

    def read_csv(self, path, index_col=0, start_date='1970-01-01', stop_date='2100-01-01', **kwargs):
        self.dataframe = pd.read_csv(path, index_col=index_col, **kwargs)
        self.dataframe.index = pd.to_datetime(self.dataframe.index)
        self.dataframe = self.dataframe[start_date:stop_date]
        
        for column in self.dataframe.columns:
            self.create_description(column)

    def find_frozen_values(
        self,
        target_column,
        window_size=3,
        var_lim_low=1,
        mode="run",
        output_column_name="preprocessed",
        outlier_column_name="frozen_values",
        return_only=False,
        **kwargs
    ):
        """
        Find frozen values.
        Creates one dataframe column called "frozen_values" where data points that are frozen
        are kept and one dataframe column called "preprocessed" where frozen data points
        are removed.
        """

        # create temporary dataframe
        frozen_values_temp_df = pd.DataFrame()

        # calc standard deviation of data within the window
        frozen_values_temp_df["Test"] = (
            self.dataframe[target_column].rolling(window_size).std(**kwargs)
        )

        # create column where 0="not frozen value", 1="frozen value"
        frozen_values_temp_df["anomalyVec"] = 0
        frozen_values_temp_df.loc[
            frozen_values_temp_df["Test"] < var_lim_low, "anomalyVec"
        ] = 1

        # calculate number of shift points if the window size is given as timestrin
        if isinstance(window_size, str):
            df_frequency = (self.dataframe.index[-1] - self.dataframe.index[0]) / (len(self.dataframe.index) - 1)
            shift_points = window_size / df_frequency

            if not shift_points.is_integer():
                window_size = round(shift_points) * df_frequency
                print('Window size not an integer of time series frequency, setting new window size to: ' + str(
                    window_size))
                shift_points = int(window_size / df_frequency)
            else:
                shift_points = int(shift_points)
        else:
            shift_points = int(window_size)

        # Process entire window as frozen
        frozen_values_temp_df["anomalyVec"] = (
            frozen_values_temp_df["anomalyVec"]
            .shift(-shift_points + 1, fill_value=0)
            .rolling(window_size)
            .max(**kwargs)
        )

        # first number of data points becomes nan due to window size, these are not detected as frozen and set to 0 here
        frozen_values_temp_df["anomalyVec"][0 : shift_points - 1] = 0

        if not return_only:
            # Create name for column
            new_column = self.create_target_column(outlier_column_name)
            new_column_2 = self.create_target_column(output_column_name)

            # Create metadata dictionary for frozen values
            metadata_dict = {"method": "find_frozen_values",
                             "used_data_column": target_column,
                             "window_size": window_size,
                             "var_lim_low": var_lim_low,
                             "mode": mode,
                             "plot_mode": "markers",
                             "plot_markers": "circle-open",
                             "marker_size": 8}
            self.create_metadata(metadata_dict, new_column)

            # Create metadata dictionary for column where frozen values have been removed from target data
            metadata_dict = {"method": "find_frozen_values",
                             "used_data_column": target_column,
                             "window_size": window_size,
                             "var_lim_low": var_lim_low,
                             "mode": mode}
            self.create_metadata(metadata_dict, new_column_2)

            # Add frozen values to new column
            self.dataframe[new_column] = (
                self.dataframe[target_column] * frozen_values_temp_df["anomalyVec"]
            ).replace(0, np.nan)

            # Add column with removed frozen values from raw data
            self.dataframe[new_column_2] = self.dataframe[target_column].loc[
                frozen_values_temp_df["anomalyVec"] < 0.5
            ]

            # Create descriptions for the new columns
            self.create_description(new_column)
            self.description[new_column]["info"] = (
                    "Data where all NON-frozen values of " + target_column + " are set to nan"
            )

            self.create_description(new_column_2)
            self.description[new_column_2]["info"] = (
                "Data where frozen values of " + target_column + " are set to nan"
            )
        else:
            frozen_return = self.dataframe[target_column] * frozen_values_temp_df["anomalyVec"].replace(0, np.nan)
            return frozen_return

    def out_of_range_detection(self, target_column, min_limit, max_limit, return_only=False):
        """
        Find values outside/inside specified limits.
        Creates one dataframe column called "out_of_range" where data points that are out of range
        are kept and one dataframe column called "preprocessed" where data points that are out of range
        are removed.
        """

        # Detect data
        temp_df = pd.DataFrame()
        temp_df["Time"] = self.dataframe.index
        temp_df = temp_df.set_index("Time")
        temp_df["out_of_range_binary"] = 0
        temp_df["out_of_range_binary"].loc[
            (self.dataframe[target_column] > max_limit)
        ] = 1
        temp_df["out_of_range_binary"].loc[
            (self.dataframe[target_column] < min_limit)
        ] = 1

        if not return_only:
            # Create name for column
            new_column = self.create_target_column("out_of_range")
            new_column_2 = self.create_target_column("preprocessed")

            # Create metadata dictionary
            metadata_dict = {"method": "out_of_range_detection",
                             "used_data_column": target_column,
                             "min_limit": min_limit,
                             "max_limit": max_limit,
                             "plot_mode": "markers",
                             "plot_markers": "circle-open",
                             "marker_size": 8}
            self.create_metadata(metadata_dict, new_column)

            # Create metadata dictionary
            metadata_dict = {"method": "out_of_range_detection",
                             "used_data_column": target_column,
                             "min_limit": min_limit,
                             "max_limit": max_limit}
            self.create_metadata(metadata_dict, new_column_2)

            self.dataframe[new_column] = (
                temp_df["out_of_range_binary"] * self.dataframe[target_column]
            ).replace(0, np.nan)

            self.dataframe[new_column_2] = self.dataframe[target_column][
                (self.dataframe[target_column] < max_limit)
                & (self.dataframe[target_column] > min_limit)
            ]

            # Create descriptions
            self.create_description(new_column)
            self.description[new_column]["info"] = (
                    "Data where all NON-out of range values of "
                    + target_column
                    + " are set to nan"
            )
            self.create_description(new_column_2)
            self.description[new_column_2]["info"] = (
                "Data where out of range values of " + target_column + " are set to nan"
            )
        else:
            out_of_range_return = (temp_df["out_of_range_binary"] * self.dataframe[target_column]).replace(0, np.nan)
            return out_of_range_return

    def write_to_excel(self, file_name):
        """
        Store data, metadata and descriptions to excel file
        """
        # create common write for all sheets
        writer = pd.ExcelWriter(file_name, engine="xlsxwriter")

        # write data to sheet "data"
        self.dataframe.to_excel(writer, sheet_name="data")

        # convert metadata dict to dataframe and write to sheet "metadata"
        md_df = pd.DataFrame.from_dict(self.metadata)
        md_df.to_excel(writer, sheet_name="metadata")

        # convert description dict to dataframe and write to sheet "description"
        dscr_df = pd.DataFrame.from_dict(self.description)
        dscr_df.to_excel(writer, sheet_name="description")

        # store file
        writer.save()
        
    def read_tsc_excel(self, file_name):
        """
        Read data, metadata and descriptions from TSC excel file
        """
        # read data
        self.dataframe = pd.read_excel(file_name, sheet_name='data', index_col=0, engine="openpyxl")

        # read metadata
        self.metadata = pd.read_excel(path, sheet_name='metadata', index_col=0, engine="openpyxl").to_dict()

        # read description
        self.description = pd.read_excel(path, sheet_name='description', index_col=0, engine="openpyxl").to_dict()
        
    def write_summary_pptx(self, presentation_name):
        """
        Creates a pptx including time series plots, metadata and descriptions of all columns in TSC dataframe.
        The user is recommended to apply format changes after exporting to pptx-format
        """
        prs = pptx.Presentation()

        # Presentation settings
        title_page_slide_layout = prs.slide_layouts[0]
        title_slide_layout = prs.slide_layouts[5]

        # Add title page
        slide = prs.slides.add_slide(title_page_slide_layout)
        title = slide.shapes.title
        title.text = "TSConstrictor Summary"
        subtitle = slide.placeholders[1]
        subtitle.text = "Autogenerated " + str(datetime.date.today())

        # Add analysis for each column in TSConstrictor dataframe
        for column in self.dataframe.columns:
            if self.dataframe[column].dtype != np.dtype('O'):
                # Create plot
                self.plot_static(y_column=column, save_name='timeseries.png')
                if not os.path.isfile('timeseries.png'):
                    continue

                # Time series slide     
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text = column + ' - time series'
                
                # Add plot to slide
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(8)
                pic = slide.shapes.add_picture("timeseries.png", left, top, width=width)
                os.remove('timeseries.png')

                # Statistics and metadata slide
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text = column + ' \n Statistics and description'
                shapes = slide.shapes

                # Find number of rows of metadata and description
                try:
                    desc_rows = len(self.description[column])
                except KeyError:
                    desc_rows = 0
                    print('No description found for ' + column)
                try:
                    meta_rows = len(self.metadata[column])
                except KeyError:
                    meta_rows = 0
                    print('No metadata found for ' + column)

                # Set number of rows to the max of metadata and description
                rows = max(meta_rows, desc_rows) + 1

                # Set column indices for metadata and description
                if desc_rows > 0 and meta_rows > 0:
                    cols = 5
                    desc_col = 0
                    meta_col = 3
                elif desc_rows + meta_rows == 0:
                    continue
                else:
                    cols = 2
                    desc_col = 0
                    meta_col = 0

                # Set table properties
                left = top = Inches(1.5)
                width = Inches(8)
                height = Inches(0.6)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # Fill table with descirption and metadata
                if desc_rows > 0:
                    # write column headings
                    table.cell(0, desc_col).text = 'Description'
                    table.cell(0, desc_col + 1).text = 'Value'

                    # write body cells
                    count = 1
                    for key in self.description[column].keys():
                        table.cell(count, desc_col).text = key
                        table.cell(count, desc_col + 1).text = str(self.description[column][key])
                        count += 1

                if meta_rows > 0:
                    # write column headings
                    table.cell(0, meta_col).text = 'Description'
                    table.cell(0, meta_col + 1).text = 'Value'

                    # write body cells
                    count = 1
                    for key in self.metadata[column].keys():
                        table.cell(count, meta_col).text = key
                        table.cell(count, meta_col + 1).text = str(self.metadata[column][key])
                        count += 1

        prs.save(presentation_name)

    def create_target_column(self, target_column):
        """
        When creating a new column, always call this function first
        like: my_target_column = self.create_target_column(my_target_column)
        It checks if it already exists a column with that name in the
        dataframe, and if it does, adds numbering (_2 / _3 / ...)
        """
        if target_column in self.dataframe.columns:
            new_column_name = target_column + "_2"
            count = 2
            while new_column_name in self.dataframe.columns:
                count = count + 1
                new_column_name = target_column + "_" + str(count)
        else:
            new_column_name = target_column
        return new_column_name

    def read_acurve_excel_export(self, excel_path):
        """
        Read excel file exported from acurve. Stores the data
        in self.dataframe and self.metadata
        """
        df = pd.read_excel(excel_path, header=None)
        df_list = []
        md_dict = dict()
        for i in range(0, len(df.columns), 3):
            temp_df = df[df.columns[i : i + 3]]
            md_dict = self.extract_acurve_excel_metadata(temp_df, md_dict)
            df_temp = self.extract_acurve_excel_data(temp_df)
            df_list.append(df_temp)
        df_merged = df_list[0]
        for df in df_list[1:]:
            df_merged = pd.merge(df_merged, df, on="Time", how="outer")
        df_merged = df_merged.sort_index()

        self.dataframe = df_merged
        self.metadata = md_dict

    def extract_acurve_excel_data(self, dataframe):
        """
        Read data of excel exports from acurve
        """
        N_ROWS_METADATA = 10
        data_name = dataframe[dataframe.columns[2]][0]
        dataframe = dataframe[N_ROWS_METADATA + 1 :]
        dataframe = dataframe.set_index(dataframe.columns[0])
        dataframe.columns = ["nan", data_name]
        dataframe = dataframe.drop(dataframe.columns[0], axis=1)
        dataframe.index = pd.to_datetime(dataframe.index, format="%Y-%M-%D %HH%MM")
        dataframe.index.name = "Time"
        return dataframe

    def extract_acurve_excel_metadata(self, dataframe, md_dict):
        """
        Read metadata of excel exports from acurve
        """
        N_ROWS_METADATA = 10
        data_name = dataframe[dataframe.columns[2]][0]
        col1 = dataframe.columns[0]
        col2 = dataframe.columns[2]
        md_dict[data_name] = dict()
        for idx, data in enumerate(dataframe[col1][0:N_ROWS_METADATA]):
            md_dict[data_name][data] = str(dataframe[col2][idx])
        return md_dict
